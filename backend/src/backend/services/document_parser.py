from __future__ import annotations

import mimetypes
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from tempfile import NamedTemporaryFile, TemporaryDirectory
from typing import Any

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from rapidocr_onnxruntime import RapidOCR

from backend.config import get_settings


MAX_EXTRACTED_TEXT = 120_000
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"}
TEXT_SUFFIXES = {".txt", ".md"}
CODE_SUFFIXES = {
    ".bat",
    ".c",
    ".cc",
    ".cpp",
    ".cs",
    ".css",
    ".go",
    ".h",
    ".hpp",
    ".html",
    ".ipynb",
    ".java",
    ".js",
    ".json",
    ".jsx",
    ".kt",
    ".m",
    ".php",
    ".py",
    ".r",
    ".rb",
    ".rs",
    ".sh",
    ".sql",
    ".swift",
    ".ts",
    ".tsx",
    ".xml",
    ".yaml",
    ".yml",
}
PRESENTATION_SUFFIXES = {".pptx", ".pptm", ".potx", ".potm"}
LEGACY_PRESENTATION_SUFFIXES = {".ppt"}


@dataclass(slots=True)
class VisualAsset:
    origin: str
    mime_type: str
    data: bytes


@dataclass(slots=True)
class ParsedDocument:
    parser_name: str
    text: str
    notes: list[str]
    images_detected: int
    visual_assets: list[VisualAsset] = field(default_factory=list)


class DocumentParser:
    """统一处理作业文件。"""

    def __init__(self) -> None:
        self._ocr_engine: RapidOCR | None = None
        self._ocr_error: str | None = None

    def parse(self, file_path: str | Path) -> ParsedDocument:
        resolved = Path(file_path).expanduser().resolve()
        suffix = resolved.suffix.lower()
        if suffix in TEXT_SUFFIXES | CODE_SUFFIXES:
            return self._trim(self._parse_text(resolved))
        if suffix == ".docx":
            return self._trim(self._parse_docx(resolved))
        if suffix == ".pdf":
            return self._trim(self._parse_pdf(resolved))
        if suffix in PRESENTATION_SUFFIXES:
            return self._trim(self._parse_presentation(resolved))
        if suffix in LEGACY_PRESENTATION_SUFFIXES:
            return self._trim(self._parse_legacy_presentation(resolved))
        if suffix in IMAGE_SUFFIXES:
            return self._trim(self._parse_image(resolved))
        raise ValueError(f"暂不支持的作业格式：{suffix}")

    def supports(self, file_path: str | Path) -> bool:
        suffix = Path(file_path).suffix.lower()
        return suffix in (
            TEXT_SUFFIXES
            | CODE_SUFFIXES
            | IMAGE_SUFFIXES
            | PRESENTATION_SUFFIXES
            | LEGACY_PRESENTATION_SUFFIXES
            | {".docx", ".pdf"}
        )

    def _get_ocr_engine(self) -> RapidOCR | None:
        if self._ocr_engine is not None:
            return self._ocr_engine
        if self._ocr_error is not None:
            return None
        try:
            self._ocr_engine = RapidOCR()
        except Exception as exc:  # pragma: no cover - 取决于本地推理环境
            self._ocr_error = str(exc)
            return None
        return self._ocr_engine

    def _ocr_image_file(self, file_path: Path) -> tuple[str, list[str]]:
        engine = self._get_ocr_engine()
        if engine is None:
            return "", [f"OCR 引擎不可用：{self._ocr_error or '未知错误'}"]

        result, _ = engine(str(file_path))
        texts: list[str] = []
        for item in result or []:
            if isinstance(item, (list, tuple)) and len(item) >= 2:
                text = item[1]
                if isinstance(text, str) and text.strip():
                    texts.append(text.strip())
        if not texts:
            return "", ["OCR 未识别出文字"]
        return "\n".join(texts), []

    def _ocr_image_bytes(self, content: bytes, suffix: str) -> tuple[str, list[str]]:
        with NamedTemporaryFile(suffix=suffix) as handle:
            handle.write(content)
            handle.flush()
            return self._ocr_image_file(Path(handle.name))

    def _guess_mime_type(self, suffix: str) -> str:
        return mimetypes.guess_type(f"file{suffix}")[0] or "image/png"

    def _convert_with_soffice(self, file_path: Path, target_format: str) -> Path | None:
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice is None:
            return None

        with TemporaryDirectory() as temp_dir:
            output_dir = Path(temp_dir)
            command = [
                soffice,
                "--headless",
                "--convert-to",
                target_format,
                "--outdir",
                str(output_dir),
                str(file_path),
            ]
            completed = subprocess.run(
                command,
                capture_output=True,
                text=True,
                check=False,
            )
            if completed.returncode != 0:
                return None

            converted = output_dir / f"{file_path.stem}.{target_format.split(':', 1)[0]}"
            if not converted.exists():
                return None
            converted_root = get_settings().artifacts_root / "converted_materials"
            converted_root.mkdir(parents=True, exist_ok=True)
            with NamedTemporaryFile(
                suffix=converted.suffix,
                dir=converted_root,
                delete=False,
            ) as handle:
                handle.write(converted.read_bytes())
                return Path(handle.name)

    def _parse_text(self, file_path: Path) -> ParsedDocument:
        for encoding in ("utf-8", "utf-8-sig", "gb18030"):
            try:
                content = file_path.read_text(encoding=encoding)
                return ParsedDocument(
                    parser_name="plain-text",
                    text=content,
                    notes=[],
                    images_detected=0,
                    visual_assets=[],
                )
            except UnicodeDecodeError:
                continue
        raise ValueError("文本文件编码无法识别。")

    def _parse_image(self, file_path: Path) -> ParsedDocument:
        image = Image.open(file_path)
        ocr_text, notes = self._ocr_image_file(file_path)
        notes = [f"图片尺寸：{image.width}x{image.height}", *notes]
        return ParsedDocument(
            parser_name="image+ocr",
            text=ocr_text,
            notes=notes,
            images_detected=1,
            visual_assets=[
                VisualAsset(
                    origin=file_path.name,
                    mime_type=self._guess_mime_type(file_path.suffix.lower()),
                    data=file_path.read_bytes(),
                )
            ],
        )

    def _parse_docx(self, file_path: Path) -> ParsedDocument:
        document = Document(str(file_path))
        notes: list[str] = []
        segments: list[str] = []

        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        if paragraphs:
            segments.append("\n".join(paragraphs))

        table_rows: list[str] = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_rows.append(" | ".join(cells))
        if table_rows:
            segments.append("[表格]\n" + "\n".join(table_rows))

        ocr_blocks: list[str] = []
        images_detected = 0
        visual_assets: list[VisualAsset] = []
        for relation in document.part.rels.values():
            if "image" not in relation.reltype:
                continue
            images_detected += 1
            suffix = Path(relation.target_ref).suffix or ".png"
            blob = relation.target_part.blob
            visual_assets.append(
                VisualAsset(
                    origin=relation.target_ref,
                    mime_type=self._guess_mime_type(suffix),
                    data=blob,
                )
            )
            text, image_notes = self._ocr_image_bytes(blob, suffix)
            notes.extend(f"文档图片：{note}" for note in image_notes)
            if text:
                ocr_blocks.append(text)

        if ocr_blocks:
            segments.append("[文档内图片 OCR]\n" + "\n\n".join(ocr_blocks))
            notes.append(f"共识别到 {images_detected} 张内嵌图片")

        return ParsedDocument(
            parser_name="docx+ocr",
            text="\n\n".join(segment for segment in segments if segment),
            notes=notes,
            images_detected=images_detected,
            visual_assets=visual_assets,
        )

    def _parse_pdf(self, file_path: Path) -> ParsedDocument:
        reader = PdfReader(str(file_path))
        notes: list[str] = []
        segments: list[str] = []
        images_detected = 0
        image_texts: list[str] = []
        visual_assets: list[VisualAsset] = []

        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                segments.append(f"[第 {page_number} 页]\n{text}")

            try:
                page_images: list[Any] = list(getattr(page, "images", []) or [])
            except Exception:
                page_images = []

            for image in page_images:
                data = getattr(image, "data", None)
                name = getattr(image, "name", f"page-{page_number}.png")
                if not data:
                    continue
                images_detected += 1
                visual_assets.append(
                    VisualAsset(
                        origin=f"page-{page_number}:{name}",
                        mime_type=self._guess_mime_type(Path(name).suffix or ".png"),
                        data=data,
                    )
                )
                text, image_notes = self._ocr_image_bytes(data, Path(name).suffix or ".png")
                notes.extend(f"PDF 图片：{note}" for note in image_notes)
                if text:
                    image_texts.append(f"[PDF 第 {page_number} 页图片 OCR]\n{text}")

        if image_texts:
            segments.append("\n\n".join(image_texts))
        if images_detected:
            notes.append(f"共识别到 {images_detected} 张 PDF 内图片")
        if not segments:
            notes.append("PDF 未提取到文字内容")

        return ParsedDocument(
            parser_name="pdf+ocr",
            text="\n\n".join(segment for segment in segments if segment),
            notes=notes,
            images_detected=images_detected,
            visual_assets=visual_assets,
        )

    def _parse_presentation(self, file_path: Path) -> ParsedDocument:
        presentation = Presentation(str(file_path))
        notes: list[str] = []
        segments: list[str] = []
        images_detected = 0
        visual_assets: list[VisualAsset] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_segments: list[str] = []

            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    slide_segments.append(text.strip())

                if getattr(shape, "has_table", False):
                    table_rows: list[str] = []
                    for row in shape.table.rows:
                        row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_values:
                            table_rows.append(" | ".join(row_values))
                    if table_rows:
                        slide_segments.append("[表格]\n" + "\n".join(table_rows))

                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    image = getattr(shape, "image", None)
                    if image is None:
                        continue
                    images_detected += 1
                    suffix = f".{getattr(image, 'ext', 'png')}"
                    blob = image.blob
                    visual_assets.append(
                        VisualAsset(
                            origin=f"slide-{slide_index}:{getattr(image, 'filename', f'image{suffix}')}",
                            mime_type=self._guess_mime_type(suffix),
                            data=blob,
                        )
                    )
                    ocr_text, image_notes = self._ocr_image_bytes(blob, suffix)
                    notes.extend(f"PPT 图片：{note}" for note in image_notes)
                    if ocr_text:
                        slide_segments.append("[图片 OCR]\n" + ocr_text)

            if slide_segments:
                segments.append(f"[第 {slide_index} 页]\n" + "\n\n".join(slide_segments))

        if images_detected:
            notes.append(f"共识别到 {images_detected} 张 PPT 内图片")
        if not segments:
            notes.append("PPT 未提取到文本内容")

        return ParsedDocument(
            parser_name="pptx+ocr",
            text="\n\n".join(segments),
            notes=notes,
            images_detected=images_detected,
            visual_assets=visual_assets,
        )

    def _parse_legacy_presentation(self, file_path: Path) -> ParsedDocument:
        converted_pdf = self._convert_with_soffice(file_path, "pdf")
        if converted_pdf is not None:
            try:
                parsed = self._parse_pdf(converted_pdf)
                return ParsedDocument(
                    parser_name="ppt-via-soffice+pdf",
                    text=parsed.text,
                    notes=["旧版 PPT 已通过 LibreOffice 转 PDF 解析", *parsed.notes],
                    images_detected=parsed.images_detected,
                    visual_assets=parsed.visual_assets,
                )
            finally:
                converted_pdf.unlink(missing_ok=True)

        raise ValueError("旧版 .ppt 需要本机安装 LibreOffice/soffice，或请先转为 .pptx 再导入。")

    def _trim(self, parsed: ParsedDocument) -> ParsedDocument:
        if len(parsed.text) <= MAX_EXTRACTED_TEXT:
            return parsed
        return ParsedDocument(
            parser_name=parsed.parser_name,
            text=parsed.text[:MAX_EXTRACTED_TEXT],
            notes=[*parsed.notes, f"文本过长，已截断到 {MAX_EXTRACTED_TEXT} 字符"],
            images_detected=parsed.images_detected,
            visual_assets=parsed.visual_assets,
        )
